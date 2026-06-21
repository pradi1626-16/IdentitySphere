"""Generate optimized hackathon presentation for IdentitySphere AI — judge-friendly, visual-first."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG   = RGBColor(0x05, 0x06, 0x0D)
RED  = RGBColor(0xE3, 0x19, 0x37)
W    = RGBColor(0xFF, 0xFF, 0xFF)
G    = RGBColor(0x8B, 0x94, 0x9E)
DK   = RGBColor(0x0D, 0x11, 0x1A)
GRN  = RGBColor(0x22, 0xC5, 0x5E)
ORG  = RGBColor(0xF9, 0x73, 0x16)
YEL  = RGBColor(0xEA, 0xB3, 0x08)
BLU  = RGBColor(0x3B, 0x82, 0xF6)

def bg(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def tx(s,l,t,w,h,txt,sz=18,c=W,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name='Calibri'; p.alignment=a
    return tb
def rl(s,l,t,w):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Pt(3)); sh.fill.solid(); sh.fill.fore_color.rgb=RED; sh.line.fill.background()
def bl(s,l,t,w,h,items,sz=13,c=G):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=f"▸  {item}"; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.name='Calibri'; p.space_before=Pt(3)
def card(s,l,t,w,h,title,val,col=RED):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=DK; sh.line.color.rgb=col; sh.line.width=Pt(1.5)
    tf=sh.text_frame; tf.word_wrap=True; tf.paragraphs[0].alignment=PP_ALIGN.CENTER
    p=tf.paragraphs[0]; p.text=val; p.font.size=Pt(26); p.font.color.rgb=col; p.font.bold=True; p.font.name='Calibri'
    p2=tf.add_paragraph(); p2.text=title; p2.font.size=Pt(10); p2.font.color.rgb=G; p2.alignment=PP_ALIGN.CENTER; p2.font.name='Calibri'
def box(s,l,t,w,h,label,col,sz=10):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=DK; sh.line.color.rgb=col; sh.line.width=Pt(2)
    tf=sh.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.text=label; p.font.size=Pt(sz); p.font.color.rgb=col; p.font.bold=True; p.alignment=PP_ALIGN.CENTER; p.font.name='Calibri'
def arrow(s,l,t,col=G):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(l),Inches(t),Inches(0.2),Inches(0.25))
    sh.fill.solid(); sh.fill.fore_color.rgb=col; sh.line.fill.background()
def nt(s,t): s.notes_slide.notes_text_frame.text=t

def build():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

    # ═══ SLIDE 1: TITLE ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rl(s,1.5,2.0,3)
    tx(s,1.5,2.2,10,1,"IdentitySphere AI",44,W,True)
    tx(s,1.5,3.3,10,0.5,"Identity Sprawl & Privileged Access Abuse Detection",22,RED,True)
    tx(s,1.5,4.1,10,0.4,"AI-Powered Cross-Platform Identity Intelligence Platform",16,G)
    tx(s,1.5,5.5,6,0.3,"Societe Generale GSC Hackathon 2026  |  Team: Pradeep M",13,G)
    tx(s,1.5,5.9,6,0.3,"Track: Identity & Access Risk Governance  |  Difficulty: Intermediate–Advanced",12,G)
    nt(s,"Welcome. IdentitySphere AI detects identity sprawl and privileged access abuse across 7 enterprise platforms using AI-driven correlation, behavioral ML, and graph-based attack path analysis. Let me walk you through what we built.")

    # ═══ SLIDE 2: PROBLEM ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,6,0.5,"The Problem",32,W,True); rl(s,0.8,0.9,2)
    tx(s,0.8,1.1,11,0.4,"Identity sprawl allows attackers to move laterally without malware",15,RED,True)
    bl(s,0.8,1.6,5.5,3,[
        "5,000+ identities scattered across AD, AWS, Azure AD, Okta, Salesforce, ServiceNow",
        "No single team sees the full cross-platform picture",
        "Over-privileged roles persist — effective privilege is hard to compute",
        "Orphaned accounts remain active months after termination",
        "Alert fatigue from noisy findings masks real abuse patterns",
    ])
    tx(s,7.2,1.1,5.5,0.4,"Real Incidents from Problem Statement",16,W,True)
    bl(s,7.2,1.6,5.5,3,[
        "Contractor AD disabled → Okta + AWS active 4 months → S3 data exfiltrated",
        "svc-etl-prod read-only in AD → inherited Global Admin via nested Azure AD group",
        "Developer API token never rotated → SaaS breach exposed production Salesforce",
        "On-call engineer temp admin on AD+AWS+Okta → privileges never revoked",
    ],12)
    nt(s,"The problem is real. Organizations manage thousands of identities across multiple platforms. Each platform has its own privilege model. When a service account gets Domain Admin in AD and S3 FullAccess in AWS, no single team sees the combined risk. These 4 real incidents from the problem statement show exactly why cross-platform visibility matters.")

    # ═══ SLIDE 3: SOLUTION ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Our Solution",32,W,True); rl(s,0.8,0.9,2)
    tx(s,0.8,1.1,10,0.4,"IdentitySphere AI — 8 core capabilities",16,RED)
    caps=[
        ("🔗","Cross-Platform\nIdentity Correlation","Maps same person across 7 platforms",BLU),
        ("🔐","Effective Privilege\nCalculator","Nested group inheritance traversal",ORG),
        ("🎯","8-Type Risk\nDetection Engine","Rule-based + ML hybrid detection",RED),
        ("🧠","Isolation Forest\nBehavioral ML","5-feature unsupervised anomaly detection",GRN),
        ("🗺️","Attack Path\nAnalysis","Graph-based lateral movement paths",ORG),
        ("💥","Blast Radius\nSimulation","What-if: revoke role → measure impact",YEL),
        ("🤖","AI Security\nCopilot","Evidence-based remediation guidance",BLU),
        ("📋","Compliance\nEngine","NIST · MITRE · GDPR · CIS auto-mapping",GRN),
    ]
    for i,(icon,title,desc,col) in enumerate(caps):
        r,c=i//4,i%4; x=0.8+c*3.1; y=1.7+r*2.7
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.8),Inches(2.3))
        sh.fill.solid(); sh.fill.fore_color.rgb=DK; sh.line.color.rgb=col; sh.line.width=Pt(1.5)
        tf=sh.text_frame; tf.word_wrap=True; tf.paragraphs[0].alignment=PP_ALIGN.CENTER
        p=tf.paragraphs[0]; p.text=icon; p.font.size=Pt(28)
        for txt,sz,tc,bd in [(title,14,col,True),(desc,10,G,False)]:
            pp=tf.add_paragraph(); pp.text=txt; pp.font.size=Pt(sz); pp.font.color.rgb=tc; pp.font.bold=bd; pp.alignment=PP_ALIGN.CENTER; pp.font.name='Calibri'
    nt(s,"Our solution has 8 core capabilities. Cross-platform identity correlation maps the same person across all 7 platforms. The effective privilege calculator traverses nested group inheritance. We run 8 types of risk detectors plus Isolation Forest ML. Attack paths are computed using NetworkX graph algorithms. The AI Copilot generates evidence-based remediation plans.")

    # ═══ SLIDE 4: ARCHITECTURE ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"System Architecture",32,W,True); rl(s,0.8,0.9,2)
    tx(s,0.8,1.1,10,0.3,"9-stage evidence pipeline — CSV → Correlation → Detection → Remediation",14,G)
    stages=[
        ("Data Sources\n(CSV/JSON)",BLU),("Identity\nResolver",ORG),("Privilege\nCalculator",YEL),
        ("Risk Detection\n(8 types)",RED),("Isolation Forest\nML",GRN),("Risk\nScoring",RED),
        ("Attack Graph\n(NetworkX)",ORG),("Blast Radius\nSimulator",YEL),("AI Copilot +\nCompliance",GRN),
    ]
    for i,(label,col) in enumerate(stages):
        x=0.5+i*1.38
        box(s,x,2.2,1.25,1.5,label,col,9)
        if i<len(stages)-1: arrow(s,x+1.27,2.8,G)
    tx(s,0.8,4.2,11,0.3,"Every finding is traceable: CSV Evidence → Correlation → Privilege → Detection → Score → Attack Path → Remediation",12,RED)
    tx(s,0.8,4.7,6,0.3,"Python · FastAPI · React · NetworkX · scikit-learn · Recharts",12,G)
    tx(s,0.8,5.0,6,0.3,"370 identities · 7 platforms · 804 accounts · 800 audit events",12,G)
    nt(s,"The architecture is a 9-stage pipeline. Raw CSV data flows through identity resolution, privilege calculation, 8 rule-based detectors plus Isolation Forest ML, composite scoring, graph-based attack path analysis, blast radius simulation, and AI-powered remediation. Every risk finding is traceable back to the underlying CSV evidence.")

    # ═══ SLIDE 5: AI/ML ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"AI & Machine Learning",32,W,True); rl(s,0.8,0.9,2)
    tx(s,0.8,1.2,6,0.4,"Isolation Forest Anomaly Detection",18,RED,True)
    bl(s,0.8,1.7,5.5,2,[
        "5 behavioral features extracted from audit events",
        "Unsupervised — no labeled training data required",
        "200 estimators, 10% contamination, scores normalized [0–100]",
        "Explainable: leave-one-out feature contributions",
    ],13)
    tx(s,0.8,3.5,6,0.3,"Hybrid Scoring",16,ORG,True)
    tx(s,0.8,3.9,6,0.3,"final_score = (rule_score × 0.6) + (ml_score × 0.4)",14,W,True)
    tx(s,7.2,1.2,5.5,0.4,"5 Behavioral Features",16,W,True)
    feats=[("login_frequency","Activity level per day"),("platform_spread","Cross-platform exposure ratio"),("privilege_to_usage","Over-provisioning indicator"),("dormancy","Days since last login"),("hour_entropy","Login timing regularity")]
    for i,(f,d) in enumerate(feats):
        y=1.7+i*0.55
        tx(s,7.2,y,2.5,0.3,f,13,RED,True)
        tx(s,9.8,y,3,0.3,d,12,G)
    tx(s,7.2,4.6,5.5,0.4,"False-Positive Suppression",16,W,True)
    bl(s,7.2,5.0,5.5,1.5,[
        "Active admin (logged in <7d) → 0.85× score",
        "MFA enabled everywhere → 0.80× score",
        "On-call personnel → 0.60× score",
        "Recent role change → 0.70× score",
    ],12)
    nt(s,"We use Isolation Forest for unsupervised behavioral anomaly detection. It trains on 5 features extracted from audit events. The hybrid formula combines rule-based scores at 60% weight with ML anomaly scores at 40%. Four suppression rules reduce false positives for legitimate high-privilege users like on-call personnel.")

    # ═══ SLIDE 6: KEY FEATURES ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Platform Features",32,W,True); rl(s,0.8,0.9,2)
    features=[
        ("Identity Inventory","370 identities with search, filter, status drill-down"),
        ("Identity Correlation","ReactFlow graph: person → accounts across 7 platforms"),
        ("Access Review","Approve / Revoke / Escalate with AI recommendations"),
        ("Privilege Explorer","User → Group → Role → Permission → Resource"),
        ("Risk Findings","505 events across 8 types with severity filtering"),
        ("Offboarding Gaps","32 terminated-but-active gaps detected"),
        ("Attack Paths","Simplified + Technical views with MITRE mapping"),
        ("Blast Radius","What-if simulation: revoke role → measure reduction"),
        ("Privilege Heatmap","Platform × Department risk concentration matrix"),
        ("AI Copilot","Evidence-based risk analysis + remediation plans"),
        ("Incident Center","Open → Review → Approve → Resolve workflow"),
        ("Scenario Simulator","5 live sim types with real-time detection"),
        ("Compliance Center","NIST · MITRE · GDPR · CIS drill-down"),
        ("5 Role Portals","Admin · Auditor · Executive · Employee · Contractor"),
    ]
    for i,(t,d) in enumerate(features):
        c=0 if i<7 else 1; r=i%7; x=0.8+c*6.2; y=1.2+r*0.82
        tx(s,x,y,2.3,0.3,t,13,W,True)
        tx(s,x+2.4,y,3.5,0.3,d,11,G)
    nt(s,"The platform has 14 admin pages and 5 role-based portals with 27 total routes. Key differentiators: the simplified attack path view lets judges understand attack chains in 5 seconds, the blast radius simulator shows risk reduction from revoking a single role, and the AI Copilot generates evidence-based remediation plans.")

    # ═══ SLIDE 7: DETECTION ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Detection Capabilities",32,W,True); rl(s,0.8,0.9,2)
    dets=[
        ("MFA Gaps","164",RED),("Privilege Escalation","110",ORG),("Cross-Platform Admins","107",RED),
        ("Orphaned Accounts","49",ORG),("Offboarding Gaps","32",YEL),("SoD Violations","31",YEL),
        ("Token Abuse","12",BLU),("Total Findings","505",RED),
    ]
    for i,(name,count,col) in enumerate(dets):
        c=i%4; r=i//4; x=0.8+c*3.1; y=1.3+r*2.5
        card(s,x,y,2.8,1.0,name,count,col)
    tx(s,0.8,4.7,11,0.3,"Alert Consolidation: 505 raw signals → 60 actionable incidents = 90.1% noise reduction",14,GRN,True)
    tx(s,0.8,5.1,11,0.3,"Exceeds problem statement target of ≥40% reduction by 2.25×",12,G)
    nt(s,"Our 8 detectors found 505 risk events. The highest volume findings are MFA gaps at 164 and privilege escalation at 110. The alert consolidation engine clusters these into 60 actionable incidents — a 90.1% reduction, exceeding the 40% target by over 2x.")

    # ═══ SLIDE 8: ATTACK PATH DEMO ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Attack Path: Live Demo Scenario",32,W,True); rl(s,0.8,0.9,2)
    tx(s,0.8,1.1,10,0.3,"How IdentitySphere detects, explains, and remediates a cross-platform attack",14,RED)
    steps=[
        ("🎯","STEP 1","Compromised\nIdentity","Manish Joshi\nOkta account","T1078",RED),
        ("🔀","STEP 2","Lateral\nMovement","Okta → AD → AWS\n5 platforms","T1550",ORG),
        ("⚡","STEP 3","Privilege\nEscalation","Domain Admin\nAdministratorAccess","T1098",YEL),
        ("💀","STEP 4","Critical\nResource","domain-controller\nFULL COMPROMISE","Impact",RED),
    ]
    for i,(icon,step,label,detail,mitre,col) in enumerate(steps):
        x=0.6+i*3.15
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.6),Inches(2.8),Inches(3.0))
        sh.fill.solid(); sh.fill.fore_color.rgb=DK; sh.line.color.rgb=col; sh.line.width=Pt(2)
        tf=sh.text_frame; tf.word_wrap=True; tf.paragraphs[0].alignment=PP_ALIGN.CENTER
        p=tf.paragraphs[0]; p.text=icon; p.font.size=Pt(36)
        for txt,sz,tc,bd in [(step,10,G,True),(label,16,col,True),(detail,12,W,False),(mitre,11,col,True)]:
            pp=tf.add_paragraph(); pp.text=txt; pp.font.size=Pt(sz); pp.font.color.rgb=tc; pp.font.bold=bd; pp.alignment=PP_ALIGN.CENTER; pp.font.name='Calibri'
        if i<3:
            arr=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+2.85),Inches(2.9),Inches(0.3),Inches(0.35))
            arr.fill.solid(); arr.fill.fore_color.rgb=col; arr.line.fill.background()
    tx(s,0.8,5.0,11,0.3,"IdentitySphere Response:",16,W,True)
    bl(s,0.8,5.4,11,1.5,[
        "Detected by rule engine + behavioral ML (anomaly score 72/100)",
        "AI Copilot: evidence-based narrative with MITRE ATT&CK mapping",
        "Remediation: Revoke AWS Admin, disable AD Domain Admin, enforce Okta MFA",
        "Blast radius reduced from 14 to 4 resources (71% reduction)",
    ],12)
    nt(s,"This is our demo scenario. An attacker compromises an Okta account, moves laterally to AD and AWS, escalates to Domain Admin, and reaches the domain controller. IdentitySphere detects this in real-time, the AI Copilot explains the risk with MITRE mapping, and generates platform-specific remediation. Blast radius drops 71% after remediation.")

    # ═══ SLIDE 9: WHY WE STAND OUT ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Why IdentitySphere AI Stands Out",32,W,True); rl(s,0.8,0.9,2)
    diffs=[
        ("Cross-Platform Correlation","Maps same person across 7 platforms — no other tool does this end-to-end",RED),
        ("Effective Privilege Calculation","Traverses nested group inheritance — reveals hidden admin access",ORG),
        ("Behavioral ML (Isolation Forest)","Unsupervised anomaly detection — no labeled data needed",GRN),
        ("Attack Path Graph Analysis","NetworkX-based lateral movement and escalation path visualization",BLU),
        ("Blast Radius Simulation","What-if: revoke one role → see exact risk reduction",YEL),
        ("AI-Powered Remediation","Evidence-based, platform-specific remediation steps",RED),
        ("90.1% Alert Reduction","505 signals → 60 incidents — 2.25× the 40% target",GRN),
        ("Full Compliance Mapping","Auto-maps to NIST 800-53 · MITRE ATT&CK · GDPR · CIS",BLU),
    ]
    for i,(t,d,col) in enumerate(diffs):
        c=0 if i<4 else 1; r=i%4; x=0.8+c*6.2; y=1.3+r*1.4
        tx(s,x,y,5.5,0.3,t,15,col,True)
        tx(s,x,y+0.3,5.5,0.5,d,12,G)
    nt(s,"What makes us stand out: end-to-end cross-platform correlation across 7 platforms, effective privilege calculation with nested group traversal, unsupervised ML that needs no labeled data, graph-based attack paths, what-if blast radius simulation, AI remediation, 90% alert reduction, and full compliance mapping. No existing tool combines all of these.")

    # ═══ SLIDE 10: RESULTS ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Results & Impact",32,W,True); rl(s,0.8,0.9,2)
    metrics=[("370","Identities\nAnalyzed",RED),("7","Platforms\nCorrelated",BLU),("505","Risk Events\nDetected",RED),("90.1%","Alert\nReduction",GRN),("100%","Identity\nCoverage",GRN)]
    for i,(v,l,col) in enumerate(metrics):
        card(s,0.6+i*2.5,1.2,2.2,1.3,l,v,col)
    tx(s,0.8,3.0,6,0.4,"All 5 Success Criteria Met ✓",20,GRN,True)
    criteria=[
        ("Identity Coverage ≥95%","100%",GRN),
        ("Risk Scenarios Identified","8 risk types",GRN),
        ("Alert Consolidation ≥40%","90.1%",GRN),
        ("Risk Explainability","5-factor breakdown",GRN),
        ("Governance Readiness","9 compliance mappings",GRN),
    ]
    for i,(label,val,col) in enumerate(criteria):
        y=3.5+i*0.45
        tx(s,0.8,y,4,0.3,label,13,G)
        tx(s,5,y,2,0.3,val,13,col,True)
    tx(s,7.5,3.0,5,0.4,"All 7 Deliverables Complete ✓",20,GRN,True)
    deliverables=["Working prototype with simulated data","Cross-platform identity resolver","Effective privilege calculator","Risk scoring with explainable breakdown","Dashboard with risk list, privileges, incidents","Architecture documentation (ARCHITECTURE.md)","Sample risk report with 10 risky identities"]
    bl(s,7.5,3.5,5.5,3,deliverables,12)
    nt(s,"All 5 success criteria from the problem statement are met. Identity coverage is 100%, alert consolidation achieves 90.1% — more than double the 40% target. All 7 deliverables are complete. The platform runs the full pipeline in about 16 seconds.")

    # ═══ SLIDE 11: COMPLIANCE ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tx(s,0.8,0.4,10,0.5,"Compliance & Framework Alignment",32,W,True); rl(s,0.8,0.9,2)
    fws=[
        ("NIST SP 800-53",["AC-2: Account Management","AC-6: Least Privilege","IA-4: Identifier Management"],BLU),
        ("MITRE ATT&CK",["T1078: Valid Accounts","T1098: Account Manipulation","T1550: Alternate Auth Material"],RED),
        ("GDPR",["Art. 5: Data Minimisation","Art. 32: Security of Processing"],GRN),
        ("CIS Controls",["Control 5: Account Management","Control 6: Access Control"],ORG),
    ]
    for i,(fw,items,col) in enumerate(fws):
        x=0.8+(i%2)*6.2; y=1.2+(i//2)*2.8
        tx(s,x,y,5.5,0.4,fw,20,col,True)
        bl(s,x,y+0.45,5.5,2,items,13)
    tx(s,0.8,6.5,11,0.3,"9 detection capabilities auto-mapped to compliance frameworks — judges can click any control to see affected identities",12,G)
    nt(s,"Every finding auto-maps to compliance frameworks. The compliance center has 9 capabilities mapped across NIST, MITRE, GDPR, and CIS Controls. This means audit teams can trace any finding directly to the relevant compliance control.")

    # ═══ SLIDE 12: CONCLUSION ═══
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rl(s,3.5,1.5,6)
    tx(s,1.5,1.7,10,0.8,"IdentitySphere AI",40,W,True,PP_ALIGN.CENTER)
    metrics2=[("370","Identities",RED),("7","Platforms",BLU),("505","Risks Detected",ORG),("90.1%","Alert Reduction",GRN)]
    for i,(v,l,col) in enumerate(metrics2):
        card(s,1.5+i*2.7,2.8,2.4,1.1,l,v,col)
    tx(s,1.5,4.5,10,0.8,'"Transforming fragmented identity data into\nactionable security intelligence through AI-driven\nrisk detection, privilege analysis, and\nexplainable remediation."',20,W,True,PP_ALIGN.CENTER)
    rl(s,4.5,5.7,4)
    tx(s,2,6.0,9,0.5,"Thank you  |  Questions?",28,RED,True,PP_ALIGN.CENTER)
    tx(s,2,6.6,9,0.3,"github.com/pradi1626-16/IdentitySphere  |  Pradeep M",13,G,False,PP_ALIGN.CENTER)
    nt(s,"To conclude: IdentitySphere AI addresses every requirement in the problem statement. 370 identities analyzed across 7 platforms, 505 risks detected with 90.1% alert reduction. All success criteria met, all deliverables complete, all compliance frameworks mapped. It transforms fragmented identity data into actionable security intelligence. Thank you — happy to take questions.")

    out="IdentitySphere_AI_Hackathon_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__=="__main__":
    build()
